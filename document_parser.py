"""
document_parser.py
PDF・PPTX・画像ファイルからテキストと画像を抽出するモジュール
"""

import io
from pathlib import Path
from PIL import Image
import pdfplumber
from pptx import Presentation
from pptx.util import Inches


def _extract_pdf_embedded_images(pdf_path: str) -> list:
    """
    PyMuPDF でPDFの各ページから埋め込み画像を抽出する（スクリーンショットではない）

    Returns:
        ページ数と同じ長さのリスト。画像がないページは None
    """
    import fitz  # pymupdf

    result = []
    doc = fitz.open(pdf_path)
    for page in doc:
        images = page.get_images(full=True)
        if not images:
            result.append(None)
            continue
        # 解像度（width×height）が最大の埋め込み画像を選ぶ
        best = max(images, key=lambda x: x[2] * x[3])
        xref = best[0]
        try:
            pix = fitz.Pixmap(doc, xref)
            if pix.n - pix.alpha > 3:          # CMYK → RGB 変換
                pix = fitz.Pixmap(fitz.csRGB, pix)
            img = Image.open(io.BytesIO(pix.tobytes("jpeg"))).convert("RGB")
            result.append(img)
        except Exception:
            result.append(None)
    doc.close()
    return result


def parse_pdf(file_path: str) -> list[dict]:
    """
    PDFからテキストと埋め込み画像を抽出する
    （埋め込み画像がないページは image=None）

    Returns:
        [{"page": int, "text": str, "image": PIL.Image|None, "source": str}, ...]
    """
    pages = []
    source = Path(file_path).name
    embedded_images = _extract_pdf_embedded_images(file_path)

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            pages.append({
                "page": i + 1,
                "text": text.strip(),
                "image": embedded_images[i] if i < len(embedded_images) else None,
                "source": source,
            })
    return pages


def parse_pptx(file_path: str) -> list[dict]:
    """
    PPTXからスライドのテキストと画像を抽出する

    Returns:
        [{"page": int, "text": str, "image": PIL.Image | None, "source": str}, ...]
    """
    prs = Presentation(file_path)
    slides = []
    source = Path(file_path).name

    for i, slide in enumerate(prs.slides):
        # テキスト抽出
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        # 画像抽出（スライド内の最初の画像を使用）
        slide_image = None
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_bytes = shape.image.blob
                slide_image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                break

        slides.append({
            "page": i + 1,
            "text": "\n".join(texts),
            "image": slide_image,
            "source": source,
        })

    return slides


def extract_site_images(url: str, max_images: int = 3) -> list:
    """
    ウェブサイトから参照用画像（OG画像・ロゴ・ヒーロー画像）をダウンロードする

    Returns:
        PIL.Image のリスト（取得できたもののみ）
    """
    import requests
    from bs4 import BeautifulSoup
    from urllib.parse import urljoin
    import io as _io

    headers = {"User-Agent": "Mozilla/5.0 (compatible; HirotecVideoAI/1.0)"}
    try:
        resp = requests.get(url, headers=headers, timeout=15)
        resp.raise_for_status()
    except Exception:
        return []

    soup = BeautifulSoup(resp.text, "html.parser")
    candidate_urls = []

    # 1. OGP画像（最優先）
    og = soup.find("meta", property="og:image")
    if og and og.get("content"):
        candidate_urls.append(urljoin(url, og["content"]))

    # 2. Twitter Card画像
    tw = soup.find("meta", attrs={"name": "twitter:image"})
    if tw and tw.get("content"):
        candidate_urls.append(urljoin(url, tw["content"]))

    # 3. ロゴ画像（src に "logo" を含む img タグ）
    for img_tag in soup.find_all("img"):
        src = img_tag.get("src", "")
        alt = img_tag.get("alt", "").lower()
        if src and ("logo" in src.lower() or "logo" in alt):
            candidate_urls.append(urljoin(url, src))
            break

    # 4. ページ内の最初の大きな画像
    for img_tag in soup.find_all("img"):
        src = img_tag.get("src", "")
        if src and not src.startswith("data:") and any(
            ext in src.lower() for ext in [".jpg", ".jpeg", ".png", ".webp"]
        ):
            candidate_urls.append(urljoin(url, src))
        if len(candidate_urls) >= max_images + 2:
            break

    # ダウンロードして PIL Image に変換
    images = []
    seen = set()
    for img_url in candidate_urls:
        if img_url in seen or len(images) >= max_images:
            break
        seen.add(img_url)
        try:
            r = requests.get(img_url, headers=headers, timeout=10)
            img = Image.open(_io.BytesIO(r.content)).convert("RGB")
            if img.width >= 100 and img.height >= 50:  # 小さすぎるアイコンは除外
                images.append(img)
        except Exception:
            continue

    return images


def parse_image(file_path: str) -> list[dict]:
    """
    画像ファイル（PNG/JPG/JPEG/WEBP）を1スライドとして読み込む

    Returns:
        [{"page": 1, "text": "", "image": PIL.Image, "source": str}]
    """
    source = Path(file_path).name
    img = Image.open(file_path).convert("RGB")
    return [{
        "page": 1,
        "text": "",
        "image": img,
        "source": source,
    }]


def load_file(file_path: str) -> list[dict]:
    """
    拡張子を見て自動的にパーサーを選択する
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        return parse_image(file_path)
    else:
        raise ValueError(f"非対応のファイル形式です: {ext}（PDF、PPTX、PNG、JPG、WEBP のみ対応）")


def parse_url(url: str) -> list[dict]:
    """
    URLのWebページをスクレイピングしてスライドリストに変換する

    ページ内の <h1>〜<h3> を区切りとして複数スライドに分割する。
    見出しがない場合は全文を1スライドにまとめる。

    Returns:
        [{"page": int, "text": str, "image": None, "source": str}, ...]
    """
    import requests
    from bs4 import BeautifulSoup

    headers = {"User-Agent": "Mozilla/5.0 (compatible; HirotecVideoAI/1.0)"}
    resp = requests.get(url, headers=headers, timeout=15)
    resp.raise_for_status()
    resp.encoding = resp.apparent_encoding

    soup = BeautifulSoup(resp.text, "html.parser")

    # 不要タグを除去
    for tag in soup(["script", "style", "nav", "footer", "header", "aside", "iframe"]):
        tag.decompose()

    source = url if len(url) <= 50 else url[:47] + "…"
    slides = []
    current_title = ""
    current_lines = []

    def flush(title, lines):
        text = (title + "\n" + "\n".join(lines)).strip()
        if text:
            slides.append({
                "page": len(slides) + 1,
                "text": text,
                "image": None,
                "source": source,
            })

    for element in soup.find_all(["h1", "h2", "h3", "p", "li", "td"]):
        tag = element.name
        text = element.get_text(separator=" ", strip=True)
        if not text:
            continue

        if tag in ("h1", "h2", "h3"):
            # 見出しが来たら前のスライドを確定
            flush(current_title, current_lines)
            current_title = text
            current_lines = []
        else:
            current_lines.append(text)
            # 1スライドが長くなりすぎたら切る（約500文字）
            if sum(len(l) for l in current_lines) >= 500:
                flush(current_title, current_lines)
                current_title = ""
                current_lines = []

    flush(current_title, current_lines)

    # 見出し区切りがなかった場合
    if not slides:
        full_text = soup.get_text(separator="\n", strip=True)[:2000]
        slides = [{"page": 1, "text": full_text, "image": None, "source": source}]

    return slides


def load_multiple_files(file_paths: list[str]) -> list[dict]:
    """
    複数ファイルを読み込み、スライドを通し番号で結合する

    Args:
        file_paths: ファイルパスのリスト

    Returns:
        全ページを通し番号で結合したリスト
        各要素: {"page": int, "text": str, "image": PIL.Image|None, "source": str}
    """
    all_pages = []
    for fp in file_paths:
        try:
            pages = load_file(fp)
            all_pages.extend(pages)
        except Exception as e:
            print(f"警告: {Path(fp).name} の読み込みをスキップしました: {e}")

    # 通し番号を振り直す
    for i, p in enumerate(all_pages):
        p["page"] = i + 1

    return all_pages
